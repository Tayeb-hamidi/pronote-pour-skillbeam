"""Parser plugins for source ingestion."""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
import hashlib
import html
import io
import json
import os
import re
from urllib.parse import parse_qs, quote_plus, urlparse

from docx import Document as DocxDocument
import httpx
from pptx import Presentation
from pypdf import PdfReader

from shared.config import get_settings
from shared.enums import SourceType

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:  # pragma: no cover - optional dependency fallback
    pdfminer_extract_text = None  # type: ignore[assignment]

try:
    from youtube_transcript_api import YouTubeTranscriptApi
except Exception:  # pragma: no cover - optional dependency fallback
    YouTubeTranscriptApi = None  # type: ignore[assignment]

try:
    from yt_dlp import YoutubeDL
except Exception:  # pragma: no cover - optional dependency fallback
    YoutubeDL = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency
    import pytesseract
except Exception:  # pragma: no cover - optional dependency fallback
    pytesseract = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency
    from PIL import Image
except Exception:  # pragma: no cover - optional dependency fallback
    Image = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency
    from pdf2image import convert_from_bytes
except Exception:  # pragma: no cover - optional dependency fallback
    convert_from_bytes = None  # type: ignore[assignment]

MAX_REMOTE_BYTES = 2 * 1024 * 1024
REMOTE_TIMEOUT_SECONDS = 12.0
REMOTE_USER_AGENT = "SkillBeamIngest/0.1"
READER_FALLBACK_PREFIX = "https://r.jina.ai/"
BOT_CHALLENGE_PATTERNS = (
    "just a moment",
    "checking your browser",
    "enable javascript and cookies",
    "attention required",
    "verify you are human",
    "captcha",
    "cloudflare",
    "datadome",
)
MAX_YOUTUBE_TRANSCRIPT_CHARS = 22000
YOUTUBE_TRANSCRIPT_LANGUAGES = ("fr", "fr-FR", "en", "en-US")
SUBTITLE_EXT_PRIORITY = ("vtt", "srv3", "ttml", "json3")


@dataclass(slots=True)
class ParsedDocument:
    """Normalized parse result."""

    text: str
    sections: list[dict]
    metadata: dict
    references: list[dict]
    source_hash: str


@dataclass(slots=True)
class RemotePayload:
    """Remote URL fetch result."""

    final_url: str
    content_type: str
    body: bytes
    truncated: bool


def compute_hash(data: bytes | str) -> str:
    """Compute SHA256 for content provenance."""

    if isinstance(data, str):
        data = data.encode("utf-8")
    return hashlib.sha256(data).hexdigest()


def parse_source(
    *,
    source_type: SourceType,
    filename: str | None,
    mime_type: str | None,
    payload_bytes: bytes | None,
    raw_text: str | None,
    link_url: str | None,
    topic: str | None,
    source_metadata: dict | None = None,
) -> ParsedDocument:
    """Parse source payload into normalized text and sections."""

    source_settings = get_settings()
    metadata = source_metadata or {}
    enable_ocr = _resolve_bool_option(
        metadata.get("enable_ocr"),
        default=source_settings.enable_ocr_fallback,
    )
    enable_table_extraction = _resolve_bool_option(
        metadata.get("enable_table_extraction"),
        default=source_settings.enable_table_extraction_default,
    )
    smart_cleaning_enabled = _resolve_bool_option(
        metadata.get("smart_cleaning"),
        default=source_settings.enable_smart_cleaning_default,
    )

    if source_type == SourceType.DOCUMENT:
        if payload_bytes is None:
            raise ValueError("missing document payload")
        text, parse_metadata = _parse_document_bytes(
            filename=filename,
            mime_type=mime_type,
            payload=payload_bytes,
            enable_ocr=enable_ocr,
            ocr_language=source_settings.ocr_language,
        )
        content_hash = compute_hash(payload_bytes)
        metadata = {"filename": filename, "mime_type": mime_type, **parse_metadata}
    elif source_type == SourceType.TEXT:
        text = (raw_text or "").strip()
        content_hash = compute_hash(text)
        metadata = {"kind": "text"}
    elif source_type == SourceType.THEME:
        metadata = source_metadata or {}
        theme_text = (topic or raw_text or "").strip()
        subject = str(metadata.get("subject") or "").strip()
        class_level = str(metadata.get("class_level") or "").strip()
        difficulty_target = str(metadata.get("difficulty_target") or "").strip()
        learning_goal = str(metadata.get("learning_goal") or "").strip()

        lines = [f"Thematique centrale: {theme_text}"]
        if subject:
            lines.append(f"Matiere: {subject}")
        if class_level:
            lines.append(f"Classe cible: {class_level}")
        if difficulty_target:
            lines.append(f"Difficulte cible: {difficulty_target}")
        if learning_goal:
            lines.append(f"Objectif pedagogique: {learning_goal}")
        lines.extend(
            [
                "",
                "Consigne de generation:",
                "- Produire un support pedagogique progressif.",
                "- Contextualiser les questions pour le niveau cible.",
                "- Favoriser la comprehension avant la memorisation.",
            ]
        )
        text = "\n".join(lines)
        content_hash = compute_hash(text)
        metadata = {
            "kind": "theme",
            "topic": theme_text,
            "subject": subject,
            "class_level": class_level,
            "difficulty_target": difficulty_target,
            "learning_goal": learning_goal,
        }
    elif source_type == SourceType.YOUTUBE:
        text, metadata = _parse_youtube_source(link_url)
        content_hash = compute_hash(text)
    elif source_type == SourceType.LINK:
        text, metadata = _parse_link_source(link_url)
        content_hash = compute_hash(text)
    elif source_type == SourceType.AUDIO_VIDEO:
        text = (
            "TODO(audio_video): pipeline transcription abstraite. "
            "Connecter un provider ASR via variables d'environnement."
        )
        content_hash = compute_hash(text)
        metadata = {"kind": "audio_video"}
    else:
        raise ValueError(f"unsupported source_type: {source_type}")

    clean_text = _normalize_whitespace(text)
    clean_text, cleaning_report = _smart_clean_text(clean_text, enabled=smart_cleaning_enabled)
    table_candidates = _extract_table_candidates(clean_text) if enable_table_extraction else []
    sections = _build_sections(clean_text)
    references = [{"source_id": "section:1", "label": "Section 1"}] if clean_text else []
    metadata["source_quality"] = _build_source_quality_report(
        text=clean_text,
        sections=sections,
        table_candidates=table_candidates,
        cleaning_report=cleaning_report,
        base_metadata=metadata,
        enable_ocr=enable_ocr,
        enable_table_extraction=enable_table_extraction,
    )
    if table_candidates:
        metadata["tables"] = table_candidates

    return ParsedDocument(
        text=clean_text,
        sections=sections,
        metadata=metadata,
        references=references,
        source_hash=content_hash,
    )


def _parse_document_bytes(
    *,
    filename: str | None,
    mime_type: str | None,
    payload: bytes,
    enable_ocr: bool,
    ocr_language: str,
) -> tuple[str, dict]:
    """Parse bytes based on extension and mime type."""

    name = (filename or "").lower()
    mime = (mime_type or "").lower()

    if name.endswith(".pdf") or mime == "application/pdf":
        return _parse_pdf(payload, enable_ocr=enable_ocr, ocr_language=ocr_language)
    if name.endswith(".docx") or mime.endswith("wordprocessingml.document"):
        return _parse_docx(payload), {"parser": "docx", "ocr_status": "not_needed"}
    if name.endswith(".pptx") or mime.endswith("presentationml.presentation"):
        return _parse_pptx(payload), {"parser": "pptx", "ocr_status": "not_needed"}
    if name.endswith((".png", ".jpg", ".jpeg")) or mime in {"image/png", "image/jpeg", "image/jpg"}:
        return _parse_image(payload, enable_ocr=enable_ocr, ocr_language=ocr_language)
    if name.endswith(".txt") or name.endswith(".md") or mime.startswith("text/"):
        return payload.decode("utf-8", errors="ignore"), {
            "parser": "text",
            "ocr_status": "not_needed",
        }

    raise ValueError(f"unsupported file format for {filename or mime_type or 'unknown'}")


def _parse_pdf(payload: bytes, *, enable_ocr: bool, ocr_language: str) -> tuple[str, dict]:
    reader = PdfReader(io.BytesIO(payload))
    page_count = len(reader.pages)
    chunks: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        chunks.append(page_text)
    extracted = "\n\n".join(chunks)
    parser = "pypdf"
    scanned_pdf_suspected = len(extracted.strip()) < max(120, page_count * 40)
    ocr_status = "not_needed"
    if len(extracted.strip()) >= 120:
        if scanned_pdf_suspected:
            ocr_status = "suggested"
        return extracted, {
            "parser": parser,
            "page_count": page_count,
            "scanned_pdf_suspected": scanned_pdf_suspected,
            "ocr_status": ocr_status,
        }

    if pdfminer_extract_text is None:
        if scanned_pdf_suspected and enable_ocr:
            ocr_text, ocr_status = _ocr_pdf_payload(payload, ocr_language=ocr_language)
            if ocr_text and len(ocr_text.strip()) > len(extracted.strip()):
                return ocr_text, {
                    "parser": "ocr_pdf",
                    "page_count": page_count,
                    "scanned_pdf_suspected": True,
                    "ocr_status": ocr_status,
                }
        return extracted, {
            "parser": parser,
            "page_count": page_count,
            "scanned_pdf_suspected": scanned_pdf_suspected,
            "ocr_status": "disabled" if scanned_pdf_suspected and not enable_ocr else "unavailable",
        }

    try:
        fallback_text = pdfminer_extract_text(io.BytesIO(payload)) or ""
        if len(fallback_text.strip()) > len(extracted.strip()):
            extracted = fallback_text
            parser = "pdfminer"
    except Exception:
        pass

    scanned_pdf_suspected = len(extracted.strip()) < max(120, page_count * 40)
    if scanned_pdf_suspected and enable_ocr:
        ocr_text, ocr_status = _ocr_pdf_payload(payload, ocr_language=ocr_language)
        if ocr_text and len(ocr_text.strip()) > len(extracted.strip()):
            return ocr_text, {
                "parser": "ocr_pdf",
                "page_count": page_count,
                "scanned_pdf_suspected": True,
                "ocr_status": ocr_status,
            }
        return extracted, {
            "parser": parser,
            "page_count": page_count,
            "scanned_pdf_suspected": True,
            "ocr_status": ocr_status,
        }

    return extracted, {
        "parser": parser,
        "page_count": page_count,
        "scanned_pdf_suspected": scanned_pdf_suspected,
        "ocr_status": "disabled" if scanned_pdf_suspected else "not_needed",
    }


def _parse_docx(payload: bytes) -> str:
    document = DocxDocument(io.BytesIO(payload))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _parse_pptx(payload: bytes) -> str:
    presentation = Presentation(io.BytesIO(payload))
    slides_text: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides_text.append(f"Slide {index}: " + " | ".join(parts))
    return "\n".join(slides_text)


def _parse_image(payload: bytes, *, enable_ocr: bool, ocr_language: str) -> tuple[str, dict]:
    if not enable_ocr:
        return "", {"parser": "image", "ocr_status": "disabled", "scanned_pdf_suspected": False}
    if Image is None or pytesseract is None:
        return "", {
            "parser": "image",
            "ocr_status": "unavailable_import",
            "scanned_pdf_suspected": False,
        }

    try:
        image = Image.open(io.BytesIO(payload))
        text = pytesseract.image_to_string(image, lang=ocr_language)  # type: ignore[arg-type]
        cleaned = _normalize_whitespace(text)
        return cleaned, {
            "parser": "image_ocr",
            "ocr_status": "applied" if cleaned else "applied_empty",
            "scanned_pdf_suspected": False,
        }
    except Exception as exc:
        return "", {
            "parser": "image",
            "ocr_status": f"failed:{_compact_error(exc)}",
            "scanned_pdf_suspected": False,
        }


def _ocr_pdf_payload(payload: bytes, *, ocr_language: str) -> tuple[str | None, str]:
    if convert_from_bytes is None or pytesseract is None:
        return None, "unavailable_import"

    try:
        pages = convert_from_bytes(payload, fmt="png", first_page=1, last_page=8)
        snippets: list[str] = []
        for page in pages:
            text = pytesseract.image_to_string(page, lang=ocr_language)  # type: ignore[arg-type]
            cleaned = _normalize_whitespace(text)
            if cleaned:
                snippets.append(cleaned)
        if not snippets:
            return None, "applied_empty"
        return "\n\n".join(snippets), "applied"
    except Exception as exc:
        return None, f"failed:{_compact_error(exc)}"


def _parse_youtube_source(link_url: str | None) -> tuple[str, dict]:
    if not link_url:
        text = "Source YouTube non fournie."
        return text, {"kind": "youtube", "url": None, "fetched": False}

    try:
        normalized_url = _normalize_url(link_url)
    except ValueError as exc:
        text = f"URL YouTube invalide: {link_url}. Detail: {exc}"
        return text, {"kind": "youtube", "url": link_url, "fetched": False, "error": str(exc)}

    video_id = _extract_youtube_video_id(normalized_url)
    oembed = _fetch_youtube_oembed(normalized_url)
    transcript_text, transcript_meta = _fetch_youtube_transcript(video_id)

    title = str((oembed or {}).get("title", "")).strip()
    author = str((oembed or {}).get("author_name", "")).strip()
    thumbnail = str((oembed or {}).get("thumbnail_url", "")).strip()

    lines = [f"Source YouTube: {normalized_url}"]
    if title:
        lines.append(f"Titre: {title}")
    if author:
        lines.append(f"Chaine: {author}")
    if video_id:
        lines.append(f"Identifiant video: {video_id}")
    if transcript_text:
        lines.append("")
        lines.append("Transcription audio (best-effort):")
        lines.append(transcript_text)
    else:
        lines.append("Transcription indisponible sur cette video.")
        lines.append("Generation basee sur les metadonnees accessibles.")

    metadata: dict[str, str | bool | None] = {
        "kind": "youtube",
        "url": normalized_url,
        "video_id": video_id,
        "fetched": bool(oembed),
        "transcript_available": bool(transcript_text),
    }
    if title:
        metadata["title"] = title
    if author:
        metadata["author_name"] = author
    if thumbnail:
        metadata["thumbnail_url"] = thumbnail
    if transcript_meta.get("transcript_error"):
        metadata["transcript_error"] = str(transcript_meta["transcript_error"])
    if transcript_meta.get("transcript_source"):
        metadata["transcript_source"] = str(transcript_meta["transcript_source"])

    return "\n".join(lines), metadata


def _parse_link_source(link_url: str | None) -> tuple[str, dict]:
    if not link_url:
        text = "URL lien non fournie."
        return text, {"kind": "link", "url": None, "fetched": False}

    try:
        normalized_url = _normalize_url(link_url)
    except ValueError as exc:
        text = f"URL invalide: {link_url}. Detail: {exc}"
        return text, {"kind": "link", "url": link_url, "fetched": False, "error": str(exc)}

    if _is_youtube_url(normalized_url):
        youtube_text, youtube_metadata = _parse_youtube_source(normalized_url)
        youtube_metadata["kind"] = "link_youtube"
        return youtube_text, youtube_metadata

    used_reader_fallback = False
    try:
        remote = _fetch_remote_payload(normalized_url)
    except Exception as exc:
        reader_remote = _fetch_reader_fallback_payload(normalized_url, cause=exc)
        if reader_remote is None:
            compact_error = _compact_remote_error(exc)
            text = (
                f"Source web: {normalized_url}\n"
                f"Recuperation impossible: {compact_error}\n"
                "Generation impossible en mode fiable pour cette URL."
            )
            return text, {
                "kind": "link",
                "url": normalized_url,
                "fetched": False,
                "error": compact_error,
            }
        remote = reader_remote
        used_reader_fallback = True

    content_type = remote.content_type.lower()
    parser = "reader" if used_reader_fallback else "text"
    title = ""
    extracted_text = ""

    if (not used_reader_fallback) and (
        "application/pdf" in content_type or remote.final_url.lower().endswith(".pdf")
    ):
        parser = "pdf"
        try:
            extracted_text = _parse_pdf(remote.body)
        except Exception as exc:
            extracted_text = f"Echec extraction PDF: {exc}"
    elif (not used_reader_fallback) and (
        "html" in content_type or remote.final_url.lower().endswith((".html", ".htm"))
    ):
        parser = "html"
        html_doc = _decode_remote_bytes(remote.body, remote.content_type)
        title, extracted_text = _extract_text_from_html(html_doc)
    else:
        extracted_text = _decode_remote_bytes(remote.body, remote.content_type)

    extracted_text = _truncate_text(_normalize_whitespace(extracted_text), max_chars=12000)
    if not extracted_text:
        extracted_text = "Aucun texte exploitable extrait de cette URL."

    if _looks_like_bot_challenge(extracted_text, title=title):
        error = "Page protegee par challenge anti-bot (Cloudflare/DataDome/CAPTCHA)"
        text = (
            f"Source web: {normalized_url}\n"
            f"Recuperation impossible: {error}\n"
            "Generation impossible en mode fiable pour cette URL."
        )
        return text, {"kind": "link", "url": normalized_url, "fetched": False, "error": error}

    lines = [f"Source web: {remote.final_url}"]
    if title:
        lines.append(f"Titre: {title}")
    lines.append("")
    lines.append(extracted_text)
    if remote.truncated:
        lines.append("")
        lines.append("[note] contenu tronque a 2MB pour securite.")

    metadata: dict[str, str | bool] = {
        "kind": "link",
        "url": remote.final_url,
        "content_type": remote.content_type or "unknown",
        "parser": parser,
        "fetched": True,
        "truncated": remote.truncated,
        "reader_fallback": used_reader_fallback,
    }
    if title:
        metadata["title"] = title

    return "\n".join(lines), metadata


def _normalize_url(raw_url: str) -> str:
    cleaned = raw_url.strip()
    if not cleaned:
        raise ValueError("URL vide")
    if "://" not in cleaned:
        cleaned = f"https://{cleaned}"
    parsed = urlparse(cleaned)
    if parsed.scheme not in {"http", "https"}:
        raise ValueError("schema non supporte")
    if not parsed.netloc:
        raise ValueError("domaine manquant")
    return parsed.geturl()


def _is_youtube_url(url: str) -> bool:
    host = urlparse(url).netloc.lower()
    return (
        host.endswith("youtube.com")
        or host.endswith("youtu.be")
        or host.endswith("youtube-nocookie.com")
    )


def _extract_youtube_video_id(url: str) -> str | None:
    parsed = urlparse(url)
    host = parsed.netloc.lower()

    if host.endswith("youtu.be"):
        slug = parsed.path.strip("/").split("/")
        return slug[0] if slug and slug[0] else None

    query = parse_qs(parsed.query)
    if query.get("v"):
        return query["v"][0]

    parts = [part for part in parsed.path.split("/") if part]
    if len(parts) >= 2 and parts[0] in {"shorts", "embed", "live", "v"}:
        return parts[1]
    return None


def _fetch_youtube_oembed(video_url: str) -> dict | None:
    oembed_url = f"https://www.youtube.com/oembed?url={quote_plus(video_url)}&format=json"
    headers = {"User-Agent": REMOTE_USER_AGENT}
    try:
        with httpx.Client(follow_redirects=True, timeout=REMOTE_TIMEOUT_SECONDS) as client:
            response = client.get(oembed_url, headers=headers)
            response.raise_for_status()
            data = response.json()
            return data if isinstance(data, dict) else None
    except Exception:
        return None


def _fetch_youtube_transcript(video_id: str | None) -> tuple[str | None, dict[str, str | bool]]:
    """Fetch YouTube transcript (captions linked to spoken audio) when available."""

    if not video_id:
        return None, {"transcript_available": False, "transcript_error": "missing_video_id"}
    if YouTubeTranscriptApi is None:
        transcript_text, transcript_meta = _fetch_youtube_transcript_with_ytdlp(video_id)
        if transcript_text:
            return transcript_text, transcript_meta
        return None, {
            "transcript_available": False,
            "transcript_error": "youtube_transcript_api_not_installed",
        }

    preferred_languages = list(YOUTUBE_TRANSCRIPT_LANGUAGES)
    last_error = ""
    proxies, cookies = _youtube_network_options()

    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(
            video_id, proxies=proxies, cookies=cookies
        )
    except Exception as exc:
        return None, {"transcript_available": False, "transcript_error": _compact_error(exc)}

    chunks: list[dict] | None = None
    source = ""

    try:
        preferred = transcript_list.find_transcript(preferred_languages)
        chunks = preferred.fetch()
        source = "captions_preferred"
    except Exception as exc:
        last_error = _compact_error(exc)

    if chunks is None:
        for transcript in transcript_list:
            if not getattr(transcript, "is_translatable", False):
                continue
            for target_language in preferred_languages:
                try:
                    chunks = transcript.translate(target_language).fetch()
                    source = f"captions_translated_{target_language}"
                    break
                except Exception:
                    continue
            if chunks is not None:
                break

    if chunks is None:
        for transcript in transcript_list:
            try:
                chunks = transcript.fetch()
                source = "captions_any_language"
                break
            except Exception as exc:
                last_error = _compact_error(exc)

    if chunks is None:
        fallback_text, fallback_meta = _fetch_youtube_transcript_with_ytdlp(video_id)
        if fallback_text:
            return fallback_text, fallback_meta
        if fallback_meta.get("transcript_error"):
            last_error = str(fallback_meta["transcript_error"])

    if chunks is None:
        return None, {
            "transcript_available": False,
            "transcript_error": last_error or "no_transcript_available",
        }

    snippets: list[str] = []
    for chunk in chunks:
        if not isinstance(chunk, dict):
            continue
        text = str(chunk.get("text", "")).replace("\n", " ").strip()
        if not text:
            continue
        if re.fullmatch(
            r"[\[\(]?\s*(music|musique|applause|rires?)\s*[\]\)]?", text, flags=re.IGNORECASE
        ):
            continue
        snippets.append(text)

    if not snippets:
        return None, {"transcript_available": False, "transcript_error": "empty_transcript"}

    transcript = _normalize_whitespace(" ".join(snippets))
    transcript = _truncate_text(transcript, max_chars=MAX_YOUTUBE_TRANSCRIPT_CHARS)
    return transcript, {"transcript_available": True, "transcript_source": source}


def _compact_error(exc: Exception) -> str:
    text = re.sub(r"\s+", " ", str(exc)).strip()
    return text[:220] if text else "unknown_error"


def _fetch_youtube_transcript_with_ytdlp(video_id: str) -> tuple[str | None, dict[str, str | bool]]:
    """Fallback transcript extraction using yt-dlp subtitle tracks."""

    if YoutubeDL is None:
        return None, {"transcript_available": False, "transcript_error": "yt_dlp_not_installed"}

    proxies, cookies = _youtube_network_options()
    watch_url = f"https://www.youtube.com/watch?v={video_id}"
    options = {
        "quiet": True,
        "no_warnings": True,
        "skip_download": True,
        "noplaylist": True,
    }
    if proxies and proxies.get("https"):
        options["proxy"] = proxies["https"]
    if cookies and os.path.exists(cookies):
        options["cookiefile"] = cookies

    try:
        with YoutubeDL(options) as ydl:
            info = ydl.extract_info(watch_url, download=False)
    except Exception as exc:
        return None, {"transcript_available": False, "transcript_error": _compact_error(exc)}

    if not isinstance(info, dict):
        return None, {"transcript_available": False, "transcript_error": "invalid_ytdlp_payload"}

    subtitle_url, subtitle_source = _pick_subtitle_track(info)
    if not subtitle_url:
        return None, {"transcript_available": False, "transcript_error": "no_ytdlp_subtitle_track"}

    try:
        response = httpx.get(
            subtitle_url,
            headers={"User-Agent": REMOTE_USER_AGENT},
            timeout=REMOTE_TIMEOUT_SECONDS,
            follow_redirects=True,
            proxy=proxies.get("https") if proxies else None,
        )
        response.raise_for_status()
    except Exception as exc:
        return None, {"transcript_available": False, "transcript_error": _compact_error(exc)}

    transcript = _subtitle_payload_to_text(response.text)
    transcript = _truncate_text(
        _normalize_whitespace(transcript), max_chars=MAX_YOUTUBE_TRANSCRIPT_CHARS
    )
    if not transcript:
        return None, {"transcript_available": False, "transcript_error": "empty_ytdlp_transcript"}

    return transcript, {"transcript_available": True, "transcript_source": subtitle_source}


def _pick_subtitle_track(info: dict) -> tuple[str | None, str]:
    subtitles = info.get("subtitles")
    automatic = info.get("automatic_captions")

    track_url = _select_subtitle_url(subtitles)
    if track_url:
        return track_url, "ytdlp_subtitles"

    track_url = _select_subtitle_url(automatic)
    if track_url:
        return track_url, "ytdlp_auto_captions"

    return None, ""


def _select_subtitle_url(tracks: object) -> str | None:
    if not isinstance(tracks, dict):
        return None

    preferred_order = list(
        dict.fromkeys([*YOUTUBE_TRANSCRIPT_LANGUAGES, "fr", "en", *tracks.keys()])
    )
    for language_code in preferred_order:
        entries = tracks.get(language_code)
        if not isinstance(entries, list):
            continue

        for ext in SUBTITLE_EXT_PRIORITY:
            for entry in entries:
                if (
                    isinstance(entry, dict)
                    and entry.get("ext") == ext
                    and isinstance(entry.get("url"), str)
                ):
                    return entry["url"]

        for entry in entries:
            if isinstance(entry, dict) and isinstance(entry.get("url"), str):
                return entry["url"]

    return None


def _subtitle_payload_to_text(payload: str) -> str:
    stripped = payload.lstrip()
    if stripped.startswith("{"):
        try:
            data = json.loads(stripped)
            if isinstance(data, dict):
                return _subtitle_json_to_text(data)
        except Exception:
            pass

    cleaned = payload.replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"(?is)<[^>]+>", " ", cleaned)

    lines: list[str] = []
    for raw_line in cleaned.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.upper().startswith("WEBVTT"):
            continue
        if line.startswith(("Kind:", "Language:", "NOTE", "STYLE", "REGION")):
            continue
        if "-->" in line:
            continue
        if re.match(r"^\d+$", line):
            continue
        lines.append(line)

    return html.unescape(" ".join(lines))


def _subtitle_json_to_text(payload: dict) -> str:
    events = payload.get("events")
    if not isinstance(events, list):
        return ""

    snippets: list[str] = []
    for event in events:
        if not isinstance(event, dict):
            continue
        segments = event.get("segs")
        if not isinstance(segments, list):
            continue
        for segment in segments:
            if not isinstance(segment, dict):
                continue
            value = str(segment.get("utf8", "")).replace("\n", " ").strip()
            if value:
                snippets.append(value)

    return html.unescape(" ".join(snippets))


def _youtube_network_options() -> tuple[dict[str, str] | None, str | None]:
    settings = get_settings()

    proxies: dict[str, str] | None = None
    if settings.youtube_proxy_url:
        proxies = {"http": settings.youtube_proxy_url, "https": settings.youtube_proxy_url}

    cookies = settings.youtube_cookies_file.strip() if settings.youtube_cookies_file else None
    if cookies and (not os.path.isfile(cookies) or os.path.getsize(cookies) == 0):
        cookies = None
    return proxies, cookies


def _fetch_remote_payload(url: str) -> RemotePayload:
    headers = {
        "User-Agent": REMOTE_USER_AGENT,
        "Accept": "text/html,text/plain,application/pdf;q=0.9,*/*;q=0.5",
    }
    with httpx.Client(follow_redirects=True, timeout=REMOTE_TIMEOUT_SECONDS) as client:
        with client.stream("GET", url, headers=headers) as response:
            response.raise_for_status()
            content_type = response.headers.get("content-type", "")
            final_url = str(response.url)
            chunks: list[bytes] = []
            total = 0
            truncated = False
            for chunk in response.iter_bytes():
                if not chunk:
                    continue
                next_total = total + len(chunk)
                if next_total > MAX_REMOTE_BYTES:
                    keep = MAX_REMOTE_BYTES - total
                    if keep > 0:
                        chunks.append(chunk[:keep])
                    truncated = True
                    break
                chunks.append(chunk)
                total = next_total

    return RemotePayload(
        final_url=final_url,
        content_type=content_type,
        body=b"".join(chunks),
        truncated=truncated,
    )


def _fetch_reader_fallback_payload(url: str, *, cause: Exception) -> RemotePayload | None:
    """Fetch web content through reader fallback when direct fetch is blocked."""

    if not _should_try_reader_fallback(cause):
        return None

    reader_url = _build_reader_fallback_url(url)
    headers = {
        "User-Agent": REMOTE_USER_AGENT,
        "Accept": "text/plain,text/markdown;q=0.9,*/*;q=0.5",
    }
    try:
        with httpx.Client(follow_redirects=True, timeout=REMOTE_TIMEOUT_SECONDS + 8) as client:
            with client.stream("GET", reader_url, headers=headers) as response:
                response.raise_for_status()
                content_type = response.headers.get("content-type", "text/plain")
                chunks: list[bytes] = []
                total = 0
                truncated = False
                for chunk in response.iter_bytes():
                    if not chunk:
                        continue
                    next_total = total + len(chunk)
                    if next_total > MAX_REMOTE_BYTES:
                        keep = MAX_REMOTE_BYTES - total
                        if keep > 0:
                            chunks.append(chunk[:keep])
                        truncated = True
                        break
                    chunks.append(chunk)
                    total = next_total
    except Exception:
        return None

    return RemotePayload(
        final_url=url,
        content_type=content_type,
        body=b"".join(chunks),
        truncated=truncated,
    )


def _should_try_reader_fallback(cause: Exception) -> bool:
    if isinstance(cause, httpx.HTTPStatusError):
        if cause.response is None:
            return False
        return 400 <= cause.response.status_code < 600
    if isinstance(cause, httpx.TransportError):
        return True
    return False


def _build_reader_fallback_url(url: str) -> str:
    return f"{READER_FALLBACK_PREFIX}{url}"


def _compact_remote_error(exc: Exception) -> str:
    if isinstance(exc, httpx.HTTPStatusError):
        status_code = exc.response.status_code if exc.response else 0
        request_url = str(exc.request.url) if exc.request else ""
        host = urlparse(request_url).netloc
        if host:
            return f"HTTP {status_code} ({host})"
        return f"HTTP {status_code}"
    return _compact_error(exc)


def _looks_like_bot_challenge(text: str, *, title: str = "") -> bool:
    haystack = f"{title}\n{text}".lower()
    return any(marker in haystack for marker in BOT_CHALLENGE_PATTERNS)


def _decode_remote_bytes(payload: bytes, content_type: str) -> str:
    charset_match = re.search(r"charset=([A-Za-z0-9._-]+)", content_type or "")
    encoding = charset_match.group(1) if charset_match else "utf-8"
    try:
        return payload.decode(encoding, errors="ignore")
    except LookupError:
        return payload.decode("utf-8", errors="ignore")


def _extract_text_from_html(html_doc: str) -> tuple[str, str]:
    title_match = re.search(r"(?is)<title[^>]*>(.*?)</title>", html_doc)
    title = ""
    if title_match:
        title = html.unescape(re.sub(r"\s+", " ", title_match.group(1))).strip()

    description_match = re.search(
        r'(?is)<meta[^>]+(?:name=["\']description["\']|property=["\']og:description["\'])[^>]*content=["\'](.*?)["\']',
        html_doc,
    )
    description = ""
    if description_match:
        description = html.unescape(re.sub(r"\s+", " ", description_match.group(1))).strip()

    cleaned = re.sub(r"(?is)<(script|style|noscript|iframe|svg).*?>.*?</\1>", " ", html_doc)
    cleaned = re.sub(r"(?is)<br\s*/?>", "\n", cleaned)
    cleaned = re.sub(r"(?is)</(p|div|li|section|article|h1|h2|h3|h4|h5|h6)>", "\n", cleaned)
    cleaned = re.sub(r"(?is)<[^>]+>", " ", cleaned)
    cleaned = html.unescape(cleaned)
    text = _normalize_whitespace(cleaned)

    if description:
        text = f"{description}\n\n{text}" if text else description

    return title, text


def _truncate_text(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rstrip() + "..."


def _resolve_bool_option(value: object, *, default: bool) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in {"1", "true", "yes", "on", "oui"}:
            return True
        if normalized in {"0", "false", "no", "off", "non"}:
            return False
    return default


def _smart_clean_text(text: str, *, enabled: bool) -> tuple[str, dict]:
    if not enabled:
        return text, {"enabled": False, "removed_lines": 0, "removed_repeated_headers": 0}

    lines = [line.strip() for line in text.splitlines()]
    normalized_counts = Counter(line.lower() for line in lines if line and len(line) <= 80)

    cleaned_lines: list[str] = []
    removed_lines = 0
    removed_repeated_headers = 0
    for line in lines:
        if not line:
            cleaned_lines.append("")
            continue

        if re.fullmatch(r"(page\s*)?\d+(\s*/\s*\d+)?", line, flags=re.IGNORECASE):
            removed_lines += 1
            continue

        if len(line) <= 80 and normalized_counts.get(line.lower(), 0) >= 4:
            removed_repeated_headers += 1
            continue

        cleaned_lines.append(line)

    compact = "\n".join(cleaned_lines)
    compact = re.sub(r"-\n([a-zà-ÿ])", r"\1", compact, flags=re.IGNORECASE)
    compact = _normalize_whitespace(compact)
    return compact, {
        "enabled": True,
        "removed_lines": removed_lines,
        "removed_repeated_headers": removed_repeated_headers,
    }


def _extract_table_candidates(text: str, max_tables: int = 8) -> list[dict]:
    candidates: list[dict] = []
    for line_index, raw_line in enumerate(text.splitlines(), start=1):
        line = raw_line.strip()
        if not line:
            continue

        cells: list[str] = []
        if line.count("|") >= 2:
            cells = [cell.strip() for cell in line.split("|") if cell.strip()]
        elif "\t" in line:
            cells = [cell.strip() for cell in line.split("\t") if cell.strip()]
        elif re.search(r"\S+\s{2,}\S+\s{2,}\S+", line):
            cells = [cell.strip() for cell in re.split(r"\s{2,}", line) if cell.strip()]

        if len(cells) < 3:
            continue
        candidates.append({"line": line_index, "columns": len(cells), "sample": cells[:5]})
        if len(candidates) >= max_tables:
            break
    return candidates


def _build_source_quality_report(
    *,
    text: str,
    sections: list[dict],
    table_candidates: list[dict],
    cleaning_report: dict,
    base_metadata: dict,
    enable_ocr: bool,
    enable_table_extraction: bool,
) -> dict:
    words = [word for word in re.split(r"\s+", text) if word]
    char_count = len(text)
    word_count = len(words)
    sentence_count = len(re.findall(r"[.!?]", text))
    readability = round((word_count / max(1, sentence_count)), 2)
    scanned_pdf_suspected = bool(base_metadata.get("scanned_pdf_suspected"))
    ocr_status = str(
        base_metadata.get("ocr_status") or ("disabled" if not enable_ocr else "unknown")
    )

    return {
        "char_count": char_count,
        "word_count": word_count,
        "sections": len(sections),
        "sentence_count": sentence_count,
        "avg_words_per_sentence": readability,
        "table_candidates": len(table_candidates),
        "table_extraction_enabled": enable_table_extraction,
        "scanned_pdf_suspected": scanned_pdf_suspected,
        "ocr_status": ocr_status,
        "smart_cleaning": cleaning_report,
    }


def _normalize_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return re.sub(r"[ \t]+", " ", text).strip()


def _build_sections(text: str) -> list[dict]:
    if not text:
        return []
    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    sections: list[dict] = []
    for index, paragraph in enumerate(paragraphs, start=1):
        sections.append({"id": f"section:{index}", "title": f"Section {index}", "text": paragraph})
    return sections
